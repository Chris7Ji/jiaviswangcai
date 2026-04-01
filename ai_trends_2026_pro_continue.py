#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
继续完成专业版PPT生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import datetime

def create_pro_ai_trends_ppt_continue():
    """继续创建专业版PPT"""
    
    # 继续第9页内容
    prs = Presentation("/Users/jiyingguo/.openclaw/workspace/2026_AI发展趋势预测_专业版_temp.pptx")
    
    # ========== 第9页：中国AI发展特色（续）==========
    # 获取第9页（索引8）
    slide = prs.slides[8]
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    
    china_points = [
        "   • 制造业数字化转型需求迫切",
        "   • 互联网应用场景全球最复杂",
        "   • 政府推动新基建与智慧城市",
        "",
        "2. 数据资源充沛",
        "   • 14亿用户产生海量数据",
        "   • 工业互联网数据快速积累",
        "   • 政府数据开放程度提升",
        "",
        "3. 政策支持有力",
        "   • 'AI+'专项行动计划",
        "   • 新基建投资超¥10万亿",
        "   • 人才培养体系完善",
        "",
        "4. 产业生态完整",
        "   • BAT等科技巨头引领",
        "   • 初创企业活跃（超5000家）",
        "   • 硬件制造供应链完善"
    ]
    
    for point in china_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.level = 1 if point.startswith("   •") else 0
    
    # ========== 第10页：风险与挑战 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "风险挑战与应对策略"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    risk_points = [
        "⚠️ 技术风险",
        "• 算法偏见与公平性问题",
        "• 数据隐私与安全挑战",
        "• 模型可解释性不足",
        "• 技术依赖与供应链风险",
        "",
        "📜 伦理与监管",
        "• 全球监管框架不统一",
        "• 就业冲击与社会影响",
        "• 军事化应用伦理争议",
        "• 知识产权保护难题",
        "",
        "💼 商业挑战",
        "• 投资回报周期较长",
        "• 人才短缺与成本上升",
        "• 技术商业化路径不清晰",
        "• 市场竞争激烈"
    ]
    
    for point in risk_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 0 if not point.startswith("•") else 1
    
    # ========== 第11页：2026年关键预测 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "2026年十大关键预测"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    predictions = [
        "1️⃣ 人形机器人进入家庭，全球销量突破100万台",
        "2️⃣ AI在制造业渗透率超45%，创造$2.1万亿价值",
        "3️⃣ 边缘AI芯片出货15亿，智能设备成为标配",
        "4️⃣ 大模型推理成本降低70%，实时交互普及",
        "5️⃣ AI辅助药物研发成功率提升50%",
        "6️⃣ 中国贡献全球40%AI专利，应用层领先",
        "7️⃣ AI生成内容占互联网流量30%",
        "8️⃣ 自动驾驶L4级在特定场景商业化",
        "9️⃣ AI伦理法规在主要经济体落地",
        "🔟 量子计算与AI结合取得突破"
    ]
    
    for pred in predictions:
        p = tf.add_paragraph()
        p.text = pred
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    
    # ========== 第12页：战略建议 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "战略建议与行动指南"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    suggestions = [
        "🎯 对企业决策者的建议",
        "• 制定AI转型路线图，分阶段实施",
        "• 建立数据治理体系，确保数据质量",
        "• 投资AI人才培养与团队建设",
        "• 探索AI+业务创新模式",
        "",
        "🔧 对技术团队的建议",
        "• 关注行业大模型，避免重复造轮子",
        "• 重视数据安全与隐私保护",
        "• 建立AI伦理审查机制",
        "• 持续学习最新技术趋势",
        "",
        "🌐 对政策制定者的建议",
        "• 完善AI法律法规体系",
        "• 加大基础研究投入",
        "• 促进产学研用协同创新",
        "• 加强国际交流合作"
    ]
    
    for point in suggestions:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 0 if not point.startswith("•") else 1
    
    # ========== 第13页：总结与展望 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "总结与展望"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle.text = (
        "AI正从'技术突破'转向'价值创造'关键阶段\n"
        "2026年将是AI规模化应用与产业深度融合的元年\n\n"
        "机遇与挑战并存，创新与责任共担\n"
        "让我们共同迎接智能新时代！\n\n"
        f"报告生成时间：{datetime.datetime.now().strftime('%Y年%m月%d日 %H:%M')}"
    )
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    
    # 保存文件
    output_path = "/Users/jiyingguo/.openclaw/workspace/2026_AI发展趋势预测_专业版.pptx"
    prs.save(output_path)
    
    return output_path

if __name__ == "__main__":
    try:
        # 先创建基础版本
        from ai_trends_2026_pro import create_pro_ai_trends_ppt
        temp_file = create_pro_ai_trends_ppt()
        
        # 然后继续完善
        output_file = create_pro_ai_trends_ppt_continue()
        print(f"✅ 专业版PPT生成成功！文件保存位置：{output_file}")
        print(f"📊 文件大小：约13页深度报告")
        print("🎨 包含内容：")
        print("   - 科技感封面设计")
        print("   - 全球投资趋势分析")
        print("   - 核心技术突破深度分析")
        print("   - 产业应用价值量化")
        print("   - 人形机器人商业化路径")
        print("   - AI+能源案例深度分析")
        print("   - 技术成熟度时间线")
        print("   - 中国AI发展特色")
        print("   - 风险挑战分析")
        print("   - 2026年十大预测")
        print("   - 战略建议与行动指南")
        print("   - 总结展望")
    except Exception as e:
        print(f"❌ PPT生成失败：{str(e)}")
        import traceback
        traceback.print_exc()