#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建精美科技感PPT - 图文并茂，每页都有设计
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import datetime

def create_beautiful_ai_ppt():
    """创建精美科技感PPT"""
    
    print("🎨 开始创建精美科技感PPT...")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置16:9宽屏
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 定义配色方案
    colors = {
        'primary': RGBColor(0, 32, 96),      # 深科技蓝
        'secondary': RGBColor(100, 200, 255), # 科技蓝
        'accent': RGBColor(255, 100, 100),    # 强调红
        'light': RGBColor(240, 245, 255),     # 浅背景
        'dark': RGBColor(10, 20, 50),         # 深背景
        'text': RGBColor(30, 30, 30),         # 正文文字
        'white': RGBColor(255, 255, 255)      # 白色
    }
    
    print("✅ 创建配色方案和模板")
    
    # ========== 第1页：科技感封面 ==========
    print("📄 创建第1页：科技感封面")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = colors['dark']
    fill.gradient_stops[1].color.rgb = colors['primary']
    
    # 添加科技感线条装饰
    for i in range(5):
        line = slide.shapes.add_shape(
            MSO_SHAPE.LINE, Inches(0), Inches(1 + i*0.5), 
            Inches(13.33), Inches(0)
        )
        line.line.color.rgb = colors['secondary']
        line.line.width = Pt(0.5)
        line.line.transparency = 0.7
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "2026年人工智能发展趋势\n深度预测与战略洞察"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = "科技驱动 · 数据洞察 · 创新引领 · 价值创造"
    p.font.size = Pt(24)
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
    tf = info_box.text_frame
    p = tf.add_paragraph()
    p.text = f"专业分析报告 · {datetime.datetime.now().strftime('%Y年%m月')} · 机密"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # ========== 第2页：精美目录 ==========
    print("📄 创建第2页：精美目录")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 浅色渐变背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = colors['light']
    fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "报告目录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 目录项（带图标和设计）
    directory_items = [
        ("📊", "全球AI投资趋势深度分析", "2022-2026详细数据与区域分布"),
        ("🚀", "核心技术突破与成熟度", "大模型、机器人、边缘AI技术进展"),
        ("🏭", "产业应用价值量化", "制造、医疗、金融领域价值创造"),
        ("🤖", "人形机器人商业化案例", "最新'通卡'机器人深度剖析"),
        ("⚡", "AI+能源创新实践", "国家能源集团项目效益分析"),
        ("🇨🇳", "中国AI发展特色优势", "四大核心竞争优势与全球定位"),
        ("⚠️", "风险挑战与应对策略", "技术、伦理、商业全方位分析"),
        ("🔮", "2026年十大关键预测", "基于数据的权威趋势预测"),
        ("🎯", "战略建议与行动指南", "企业、技术、政策层面建议")
    ]
    
    # 创建两列目录
    left_x, right_x = Inches(1), Inches(7)
    start_y = Inches(1.8)
    item_height = Inches(0.8)
    
    for i, (icon, title, desc) in enumerate(directory_items):
        x_pos = left_x if i % 2 == 0 else right_x
        y_pos = start_y + (i // 2) * item_height
        
        # 目录项容器
        item_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(5), Inches(0.7))
        tf = item_box.text_frame
        
        # 图标和标题
        p = tf.add_paragraph()
        p.text = f"{icon} {title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        
        # 描述
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = colors['text']
        p.space_before = Pt(2)
    
    # ========== 第3页：全球投资趋势（带设计）==========
    print("📄 创建第3页：全球投资趋势（精美设计）")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['light']
    
    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "📊 全球AI投资趋势深度分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 左侧：投资数据卡片
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(5.5), Inches(2)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = colors['white']
    card1.line.color.rgb = colors['secondary']
    card1.line.width = Pt(1)
    
    card1_text = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(5), Inches(1.8))
    tf = card1_text.text_frame
    tf.clear()
    
    investment_data = [
        "📈 2022-2026投资增长轨迹",
        "• 2022年：$3200亿（基础研究期）",
        "• 2023年：$4500亿（技术验证期）",
        "• 2024年：$5800亿（商业化探索）",
        "• 2025年：$6500亿（规模应用启动）",
        "• 2026预测：$7200亿（价值创造深化）"
    ]
    
    for item in investment_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # 右侧：区域分布卡片
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.5), Inches(2)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = colors['white']
    card2.line.color.rgb = colors['secondary']
    card2.line.width = Pt(1)
    
    card2_text = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5), Inches(1.8))
    tf = card2_text.text_frame
    tf.clear()
    
    region_data = [
        "🌍 区域分布格局",
        "• 北美：42%（技术研发领先）",
        "• 亚洲：38%（应用创新活跃）",
        "• 欧洲：15%（伦理监管严格）",
        "• 其他：5%（新兴市场增长）"
    ]
    
    for item in region_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # 底部：投资主体卡片
    card3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.5)
    )
    card3.fill.solid()
    card3.fill.fore_color.rgb = colors['white']
    card3.line.color.rgb = colors['secondary']
    card3.line.width = Pt(1)
    
    card3_text = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(11.3), Inches(1.3))
    tf = card3_text.text_frame
    tf.clear()
    
    investor_data = [
        "💼 投资主体构成",
        "科技巨头：45% | 风险投资：25% | 企业投资：20% | 政府投资：10%"
    ]
    
    for item in investor_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
    
    # ========== 第4页：核心技术突破（精美设计）==========
    print("📄 创建第4页：核心技术突破（精美设计）")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['light']
    
    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "🚀 核心技术突破与成熟度"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 大语言模型技术卡片
    llm_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(5.5), Inches(3)
    )
    llm_card.fill.solid()
    llm_card.fill.fore_color.rgb = colors['white']
    llm_card.line.color.rgb = RGBColor(100, 200, 255)
    llm_card.line.width = Pt(2)
    
    llm_text = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(5), Inches(2.8))
    tf = llm_text.text_frame
    tf.clear()
    
    llm_content = [
        "🧠 大语言模型突破",
        "• 参数规模：万亿→十万亿级",
        "• 推理成本：降低70%，实时交互普及",
        "• 多模态能力：文本+图像+视频+3D统一",
        "• 行业专业化：金融、医疗垂直大模型成熟"
    ]
    
    for item in llm_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # 机器人技术卡片
    robot_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.5), Inches(3)
    )
    robot_card.fill.solid()
    robot_card.fill.fore_color.rgb = colors['white']
    robot_card.line.color.rgb = RGBColor(255, 150, 100)
    robot_card.line.width = Pt(2)
    
    robot_text = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5), Inches(2.8))
    tf = robot_text.text_frame
    tf.clear()
    
    robot_content = [
        "🤖 机器人技术商业化",
        "• 人形机器人：成本$50万→$20万",
        "• 运动控制：双足行走稳定性突破",
        "• 感知交互：多传感器自然交互",
        "• 自主决策：复杂环境适应能力"
    ]
    
    for item in robot_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_