#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建最终精美版PPT - 每页都有科技感设计，图文并茂
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import datetime

def create_beautiful_ppt():
    """创建精美科技感PPT"""
    
    print("🎨 开始创建精美科技感PPT...")
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 配色方案
    colors = {
        'primary': RGBColor(0, 32, 96),      # 深科技蓝
        'secondary': RGBColor(100, 200, 255), # 科技蓝
        'accent': RGBColor(255, 100, 100),    # 强调红
        'light': RGBColor(240, 245, 255),     # 浅背景
        'text': RGBColor(30, 30, 30),         # 正文文字
        'white': RGBColor(255, 255, 255)      # 白色
    }
    
    # ========== 第1页：科技感封面 ==========
    print("📄 创建第1页：科技感封面")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(10, 20, 50)
    fill.gradient_stops[1].color.rgb = colors['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "2026年人工智能发展趋势\n精美科技感深度报告"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = "图文并茂 · 数据驱动 · 设计精美 · 深度洞察"
    p.font.size = Pt(24)
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
    tf = info_box.text_frame
    p = tf.add_paragraph()
    p.text = f"专业分析报告 · {datetime.datetime.now().strftime('%Y年%m月%d日')} · 机密"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # ========== 第2页：精美目录 ==========
    print("📄 创建第2页：精美目录")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 浅色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['light']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "📋 报告目录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 目录项（3x3网格卡片设计）
    items = [
        ("📊", "全球AI投资趋势", "2022-2026详细数据分析"),
        ("🚀", "核心技术突破", "大模型、机器人技术进展"),
        ("🏭", "产业应用价值", "制造、医疗、金融量化分析"),
        ("🤖", "人形机器人案例", "最新'通卡'机器人深度"),
        ("⚡", "AI+能源实践", "国家能源集团项目剖析"),
        ("🇨🇳", "中国AI特色", "四大优势与全球定位"),
        ("⚠️", "风险挑战", "技术、伦理、商业分析"),
        ("🔮", "2026年预测", "十大关键趋势预测"),
        ("🎯", "战略建议", "企业、技术、政策指南")
    ]
    
    # 创建3x3网格
    for i, (icon, title, desc) in enumerate(items):
        row = i // 3
        col = i % 3
        
        x = Inches(0.8 + col * 4.2)
        y = Inches(1.8 + row * 1.8)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = colors['white']
        card.line.color.rgb = colors['secondary']
        card.line.width = Pt(1)
        
        # 卡片内容
        text_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(1.2))
        tf = text_box.text_frame
        tf.clear()
        
        # 图标和标题
        p = tf.add_paragraph()
        p.text = f"{icon} {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        
        # 描述
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = colors['text']
        p.space_before = Pt(4)
    
    # ========== 第3页：全球投资趋势（精美设计）==========
    print("📄 创建第3页：全球投资趋势（精美设计）")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['light']
    
    # 标题区域
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "📊 全球AI投资趋势深度分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 左侧卡片：投资增长
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(5.5), Inches(2.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = colors['white']
    card1.line.color.rgb = colors['secondary']
    card1.line.width = Pt(2)
    
    card1_text = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(5), Inches(2.3))
    tf = card1_text.text_frame
    tf.clear()
    
    growth_data = [
        "📈 2022-2026投资增长",
        "2022: $3200亿 | 2023: $4500亿",
        "2024: $5800亿 | 2025: $6500亿",
        "2026预测: $7200亿",
        "年复合增长率: 28%"
    ]
    
    for item in growth_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("202") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(6)
    
    # 右侧卡片：区域分布
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.5), Inches(2.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = colors['white']
    card2.line.color.rgb = colors['accent']
    card2.line.width = Pt(2)
    
    card2_text = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5), Inches(2.3))
    tf = card2_text.text_frame
    tf.clear()
    
    region_data = [
        "🌍 区域分布格局",
        "北美: 42%（技术研发领先）",
        "亚洲: 38%（应用创新活跃）",
        "欧洲: 15%（伦理监管严格）",
        "其他: 5%（新兴市场增长）"
    ]
    
    for item in region_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("北") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(6)
    
    # 底部卡片：投资主体
    card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5), Inches(11.7), Inches(1.5))
    card3.fill.solid()
    card3.fill.fore_color.rgb = colors['white']
    card3.line.color.rgb = RGBColor(100, 200, 100)
    card3.line.width = Pt(2)
    
    card3_text = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.3), Inches(1.1))
    tf = card3_text.text_frame
    tf.clear()
    
    investor_data = [
        "💼 投资主体构成",
        "科技巨头: 45% | 风险投资: 25% | 企业投资: 20% | 政府投资: 10%"
    ]
    
    for item in investor_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
    
    # ========== 第4页：核心技术突破（精美设计）==========
    print("📄 创建第4页：核心技术突破（精美设计）")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['light']
    
    # 标题区域
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "🚀 核心技术突破与成熟度"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 技术卡片1：大语言模型
    tech1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(5.5), Inches(3))
    tech1.fill.solid()
    tech1.fill.fore_color.rgb = colors['white']
    tech1.line.color.rgb = RGBColor(100, 200, 255)
    tech1.line.width = Pt(2)
    
    tech1_text = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(5), Inches(2.8))
    tf = tech1_text.text_frame
    tf.clear()
    
    llm_content = [
        "🧠 大语言模型突破",
        "• 参数规模: 万亿→十万亿级",
        "• 推理成本: 降低70%，实时交互",
        "• 多模态: 文本+图像+视频+3D",
        "• 专业化: 行业垂直大模型成熟"
    ]
    
    for item in llm_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # 技术卡片2：机器人技术
    tech2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.5), Inches(3))
    tech2.fill.solid()
    tech2.fill.fore_color.rgb = colors['white']
    tech2.line.color.rgb = RGBColor(255, 150, 100)
    tech2.line.width = Pt(2)
    
    tech2_text = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5), Inches(2.8))
    tf = tech2_text.text_frame
    tf.clear()
    
    robot_content = [
        "🤖 机器人技术商业化",
        "• 人形机器人: $50万→$20万",
        "• 运动控制: 双足行走稳定",
        "• 感知交互: 多传感器融合",
        "• 自主决策: 复杂环境适应"
    ]
    
    for item in robot_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # 技术卡片3：边缘AI
    tech3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
    tech3.fill.solid()
    tech3.fill.fore_color.rgb = colors['white']
    tech3.line.color.rgb = RGBColor(150, 100, 255)
    tech3.line.width = Pt(2)
    
    tech3_text = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(11.3), Inches(1.1))
    tf = tech3_text.text_frame
    tf.clear()
    
    edge_content = [
        "⚡ 边缘AI与AI芯片",
        "• 边缘AI芯片出货: 15亿（2026预测）",
        "• 端侧设备智能化: 成为标配",
        "• 功耗降低: 50%，性能提升3倍"
    ]
    
    for item in edge_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("•") else Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # ========== 第5页：产业应用（精美设计）==========
    print("📄 创建第5页：产业应用（精美设计）")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['light']
    
    # 标题区域
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "🏭 产业应用价值量化分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    #