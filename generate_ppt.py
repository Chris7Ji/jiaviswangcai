#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenClaw 使用指南 PPT 生成器
创建精美的PPT演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色方案 - 科技蓝主题
COLORS = {
    'primary': RGBColor(0x1E, 0x3A, 0x5F),      # 深蓝
    'secondary': RGBColor(0x2E, 0x5C, 0x8A),    # 中蓝
    'accent': RGBColor(0xF4, 0x7C, 0x20),       # 橙色点缀
    'white': RGBColor(0xFF, 0xFF, 0xFF),        # 白色
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),   # 浅灰
    'dark_text': RGBColor(0x33, 0x33, 0x33),    # 深灰文字
}

def add_title_slide(prs, title, subtitle, speaker):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                   prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # 添加演讲者信息
    speaker_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.8))
    tf = speaker_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"演讲者: {speaker}"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, section_title):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                   prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['secondary']
    shape.line.fill.background()
    
    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"第{section_num}部分"
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, has_subtitle=False, subtitle=""):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                 prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()
    
    # 标题背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                      prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 副标题（如果有）
    start_y = 1.6
    if has_subtitle and subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['secondary']
        p.font.italic = True
        start_y = 1.8
    
    # 内容列表
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 检查是否是子项（以空格开头）
        if item.startswith("  ") or item.startswith("   "):
            p.text = "    • " + item.strip()
            p.font.size = Pt(16)
            p.level = 1
        else:
            p.text = "• " + item
            p.font.size = Pt(20)
            p.level = 0
        
        p.font.color.rgb = COLORS['dark_text']
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """添加双栏对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                 prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()
    
    # 标题背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                      prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_text']
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_text']
        p.space_after = Pt(8)
    
    return slide

def add_code_slide(prs, title, code_lines):
    """添加代码展示页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                 prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()
    
    # 标题背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                      prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 代码背景
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), 
                                     Inches(12.333), Inches(5.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = COLORS['light_gray']
    code_bg.line.color.rgb = COLORS['secondary']
    
    # 代码内容
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.933), Inches(5.1))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    code_text = "\n".join(code_lines)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = COLORS['dark_text']
    
    return slide

def add_end_slide(prs, title, contact):
    """添加结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                   prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 联系信息
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = contact
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== 开始创建PPT ====================

print("🎨 开始创建 OpenClaw 使用指南 PPT...")

# 第1页：封面
add_title_slide(prs, "OpenClaw 使用指南", "打造你的AI智能助理", "季")
print("✅ 封面页")

# 第2页：什么是OpenClaw
add_content_slide(prs, "什么是OpenClaw？", [
    "OpenClaw 是AI智能体平台，不只是聊天对话",
    "核心能力：理解意图 → 执行任务",
    "与传统AI的区别：",
    "  ChatGPT/Claude: 聊天对话，回答问题",
    "  OpenClaw: 真正执行任务，自动化工作",
    "核心架构：Gateway（网关）+ Agent（智能体）+ Skills（技能）"
])
print("✅ OpenClaw介绍")

# 第3页：核心架构
add_two_column_slide(prs, "OpenClaw 核心架构", 
    "三层架构",
    [
        "Gateway（网关）",
        "  消息路由与分发",
        "  连接各通信渠道",
        "Agent（智能体）",
        "  理解任务意图",
        "  调用Skills执行",
        "Skills（技能）",
        "  具体功能实现",
        "  可扩展的工具集"
    ],
    "类比理解",
    [
        "大脑：大语言模型",
        "  负责理解和推理",
        "手脚：Skills工具",
        "  执行具体操作",
        "神经系统：Gateway",
        "  传递信息和指令",
        "记忆：Memory系统",
        "  存储经验和偏好"
    ]
)
print("✅ 核心架构")

# 第4页：安装前准备
add_content_slide(prs, "安装前准备", [
    "系统要求：",
    "  macOS / Linux / Windows(WSL)",
    "  Node.js 18+",
    "  稳定的网络连接",
    "必要账号：",
    "  大模型API Key（Kimi/Deepseek等）",
    "  飞书/钉钉等（可选，用于消息接收）",
    "准备工作：",
    "  确保系统满足要求",
    "  准备好API Key",
    "  预留足够的磁盘空间"
])
print("✅ 安装准备")

# 第5页：安装步骤（代码页）
add_code_slide(prs, "安装步骤", [
    "# 1. 安装 OpenClaw CLI",
    "npm install -g openclaw",
    "",
    "# 2. 初始化配置",
    "openclaw init",
    "",
    "# 3. 配置大模型",
    "openclaw configure",
    "",
    "# 4. 验证安装",
    "openclaw status"
])
print("✅ 安装步骤")

# 第6页：核心Skills
add_content_slide(prs, "必装核心 Skills", [
    "Clawsec - 安全防护",
    "  安装Skills前自动安全扫描",
    "Tavily Search - AI实时搜索",
    "  解决大模型知识截止问题",
    "Multi Search Engine - 多引擎聚合",
    "  17个搜索引擎，智能分流",
    "Self-Improving Agent - 自我进化",
    "  记录错误，持续学习优化",
    "GitHub - 代码管理",
    "  自然语言操作GitHub",
    "Office-Automation - 办公助手",
    "  日程、邮件、文档自动化"
])
print("✅ 核心Skills")

# 第7页：章节分隔 - 第二部分
add_section_slide(prs, "二", "记忆系统与Skills深度解析")
print("✅ 第二部分分隔")

# 第8页：记忆系统架构
add_two_column_slide(prs, "记忆系统架构",
    "三层记忆体系",
    [
        "会话级记忆",
        "  当前对话上下文",
        "  临时工作数据",
        "项目级记忆",
        "  MEMORY.md 核心记忆",
        "  .learnings/ 学习记录",
        "技能级记忆",
        "  SKILL.md 技能说明",
        "  AGENTS.md 配置规则"
    ],
    "核心记忆文件",
    [
        "SOUL.md",
        "  身份、性格、边界",
        "MEMORY.md",
        "  核心决策和偏好",
        "AGENTS.md",
        "  Agent团队配置",
        "USER.md",
        "  老板信息和偏好",
        "TOOLS.md",
        "  工具配置和环境"
    ]
)
print("✅ 记忆系统架构")

# 第9页：记忆系统实战
add_content_slide(prs, "记忆系统实战：如何记录学习", [
    "学习记录 LEARNINGS.md",
    "  成功经验、最佳实践",
    "  用户纠正、更好的方法",
    "错误记录 ERRORS.md",
    "  命令失败、API错误",
    "  知识过时、工具异常",
    "功能请求 FEATURE_REQUESTS.md",
    "  缺失的能力",
    "  新功能需求",
    "自动归档策略：每周回顾、每月归档、每季审查"
])
print("✅ 记忆系统实战")

# 第10页：核心Skills介绍1
add_content_slide(prs, "核心 Skills：信息获取类", [
    "Tavily Search - AI优化搜索",
    "  结构化输出，不是网页链接",
    "  AI优化，减少文本处理负担",
    "  上下文感知，优化搜索关键词",
    "Multi Search Engine - 多引擎聚合",
    "  17个搜索引擎（8国内+9国际）",
    "  无需API Key，智能分流",
    "  结果去重，综合呈现",
    "使用建议：",
    "  实时新闻 → Tavily",
    "  综合研究 → Multi Search"
])
print("✅ Skills介绍1")

# 第11页：核心Skills介绍2
add_content_slide(prs, "核心 Skills：自我进化类", [
    "Self-Improving Agent - 持续学习",
    "  自动监控命令执行结果",
    "  结构化记录到.learnings/",
    "  智能检索历史记录",
    "Proactive Agent - 主动服务",
    "  心跳机制：每15分钟自动唤醒",
    "  任务监控：持续跟踪进行中的任务",
    "  自我迭代：优化工作流程",
    "长期价值：",
    "  从被动执行到主动服务",
    "  越用越懂你的工作习惯"
])
print("✅ Skills介绍2")

# 第12页：Skills安装与管理
add_code_slide(prs, "Skills 安装与管理", [
    "# 安装 Skills",
    "clawhub install <skill-name>",
    "",
    "# 查看已安装",
    "clawhub list",
    "",
    "# 更新 Skills",
    "clawhub update",
    "",
    "# 搜索 Skills",
    "clawhub search <keyword>",
    "",
    "# 查看技能信息",
    "clawhub info <skill-name>"
])
print("✅ Skills管理")

# 第13页：章节分隔 - 第三部分
add_section_slide(prs, "三", "实战案例演示")
print("✅ 第三部分分隔")

# 第14页：案例1 - 定时任务
add_content_slide(prs, "案例1：定时任务设置", [
    "场景：每天早上自动获取AI新闻",
    "配置步骤：",
    "  1. 创建定时任务（cron表达式）",
    "  2. 设置执行时间：06:30",
    "  3. 配置任务内容：搜索AI新闻",
    "  4. 设置通知方式：飞书+邮件",
    "实际运行的定时任务：",
    "  • AI新闻每日摘要 - 06:30",
    "  • OpenClaw新闻监控 - 06:00",
    "  • 健康长寿科研成果 - 08:30"
])
print("✅ 案例1定时任务")

# 第15页：案例2 - 自动发送邮件
add_code_slide(prs, "案例2：自动发送邮件", [
    "# 邮件配置",
    "SMTP服务器: smtp.qq.com:465",
    "发件邮箱: 86940135@qq.com",
    "授权码: swqfjvmoupdebhgh",
    "收件人: 86940135@qq.com, jiyingguo@huawei.com",
    "",
    "# 自动发送脚本",
    "bash auto_send_email.sh",
    "",
    "# 发送内容",
    "• AI新闻每日摘要",
    "• 健康长寿科研报告",
    "• OpenClaw技术动态"
])
print("✅ 案例2邮件发送")

# 第16页：案例3 - 语音交互
add_content_slide(prs, "案例3：语音交互", [
    "TTS语音合成",
    "  文本转语音，自动朗读",
    "  支持多种音色选择",
    "  可调节语速和音调",
    "语音发送飞书",
    "  生成语音文件",
    "  自动发送到飞书",
    "  支持长文本分段",
    "配置说明：",
    "  Azure TTS服务端点",
    "  首选音色：zh-CN-XiaoyiNeural",
    "  语速：+30%（比正常快30%）"
])
print("✅ 案例3语音交互")

# 第17页：案例4 - 撰写报告
add_content_slide(prs, "案例4：自动撰写报告", [
    "工作流程：",
    "  1. 搜索相关资料（Tavily/Multi Search）",
    "  2. 整理内容结构（AI分析）",
    "  3. 生成Markdown格式",
    "  4. 转换为HTML/Word/PDF",
    "实际案例：",
    "  • OpenClaw技能总结文章",
    "  • 伊朗战争局势分析",
    "  • 自动表格生成",
    "  • 微信公众号文章撰写",
    "输出格式：",
    "  Markdown、HTML、表格、PPT"
])
print("✅ 案例4撰写报告")

# 第18页：案例5 - Agent团队协作
add_two_column_slide(prs, "案例5：Agent团队协作",
    "Agent团队",
    [
        "总指挥(main)",
        "  任务识别与分发",
        "笔杆子(creator)",
        "  内容创作与编辑",
        "参谋(canmou)",
        "  深度研究与分析",
        "进化官(jinhua)",
        "  代码开发与优化",
        "交易官(jiaoyi)",
        "  股票监控与分析",
        "社区官(shequ)",
        "  社区运营与管理"
    ],
    "协作流程",
    [
        "1. 任务识别",
        "  分析用户需求",
        "  确定任务类型",
        "2. Agent派发",
        "  选择合适Agent",
        "  分配子任务",
        "3. 并行执行",
        "  多个Agent同时工作",
        "  提高效率",
        "4. 结果汇总",
        "  整合各Agent输出",
        "  统一汇报给老板"
    ]
)
print("✅ 案例5Agent协作")

# 第19页：案例6 - 一键生成公众号文章并发布 ⭐NEW
add_content_slide(prs, "案例6：一键生成公众号文章并发布 ⭐", [
    "场景：多Agent协作撰写微信公众号文章",
    "工作流程：",
    "  1. 老板提出需求（主题、字数、风格）",
    "  2. 参谋Agent深度调研（OpenClaw前世今生）",
    "  3. 笔杆子Agent撰写文章（2800字高质量文章）",
    "  4. 尝试API自动发布（遇到个人订阅号限制）",
    "  5. 生成手工发布辅助工具（优化格式+详细指南）",
    "成果：",
    "  • 调研报告：GitHub 297k Stars等核心数据",
    "  • 高质量文章：5个部分，故事性强，金句频出",
    "  • 发布方案：API尝试→第三方工具→手工发布辅助",
    "耗时：多Agent协作仅4分19秒"
])
print("✅ 案例6公众号文章生成")

# 第20页：高级技巧
add_content_slide(prs, "高级技巧与优化", [
    "记忆检索优化：",
    "  创建关键词索引（INDEX.md）",
    "  分类存储学习记录",
    "  禁用慢速memory_search",
    "响应速度优化：",
    "  预加载核心记忆文件",
    "  缓存常用数据",
    "  使用快速路径",
    "故障排查：",
    "  查看日志定位问题",
    "  记录错误到ERRORS.md",
    "  自动修复常见问题"
])
print("✅ 高级技巧")

# 第21页：价值总结
add_content_slide(prs, "OpenClaw 价值总结", [
    "核心价值：",
    "  从聊天到执行 - 真正完成任务",
    "  从被动到主动 - 主动推送信息",
    "  从通用到个性化 - 越用越懂你",
    "适用场景：",
    "  信息收集与整理 - 新闻监控、研究报告",
    "  内容创作与编辑 - 文章撰写、报告生成",
    "  日常办公自动化 - 邮件、日程、提醒",
    "  技术开发辅助 - 代码、调试、GitHub",
    "最终目标：",
    "  成为你的AI智能助理，提升工作效率"
])
print("✅ 价值总结")

# 第22页：结束页
add_end_slide(prs, "谢谢聆听！", "Q&A 欢迎交流")
print("✅ 结束页")

# 保存PPT
output_path = "/Users/jiyingguo/.openclaw/workspace/OpenClaw使用指南.pptx"
prs.save(output_path)
print(f"\n🎉 PPT生成完成！")
print(f"📁 保存路径: {output_path}")
print(f"📊 总页数: {len(prs.slides)} 页")