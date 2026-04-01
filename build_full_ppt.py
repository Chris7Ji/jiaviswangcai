import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

print("Starting full PPT generation...")

prs = Presentation()

# Colors
MAIN_TERRACOTTA = RGBColor(165, 129, 105) # #A58169
MAIN_SAGE = RGBColor(111, 122, 97) # #6F7A61
BG_MORANDI_GRAY = RGBColor(240, 235, 229) # #F0EBE5
CONTENT_WARM_WHITE = RGBColor(253, 251, 246) # #FDFBF6
ACCENT_ROSE = RGBColor(201, 160, 155) # #C9A09B
TEXT_DARK = RGBColor(80, 80, 80)

def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_MORANDI_GRAY

def apply_font(paragraph, size=18, bold=False, color=TEXT_DARK, name='Microsoft YaHei', align=PP_ALIGN.LEFT):
    paragraph.font.name = name
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.alignment = align
    # Set line spacing to 1.5 lines
    paragraph.line_spacing = 1.5

def create_cover_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title slide
    set_bg(slide)
    
    title_box = slide.shapes.title
    title_box.text = title
    for p in title_box.text_frame.paragraphs:
        apply_font(p, size=40, bold=True, color=MAIN_TERRACOTTA, align=PP_ALIGN.CENTER)
        
    sub_box = slide.placeholders[1]
    sub_box.text = subtitle
    for p in sub_box.text_frame.paragraphs:
        apply_font(p, size=24, color=MAIN_SAGE, align=PP_ALIGN.CENTER)

def create_content_slide(title, content_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    set_bg(slide)
    
    # Add Warm White Content Area (70% width to keep generous margins)
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(5.5)
    rect = slide.shapes.add_shape(1, left, top, width, height) # 1 = rectangle
    rect.fill.solid()
    rect.fill.fore_color.rgb = CONTENT_WARM_WHITE
    rect.line.fill.background()
    
    # Title Box
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(8.0), Inches(1.0))
    tf_title = tb.text_frame
    tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_title.paragraphs[0]
    p.text = title.strip()
    apply_font(p, size=32, bold=True, color=MAIN_SAGE, align=PP_ALIGN.LEFT)
    
    # Content Box
    cb = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7.0), Inches(4.8))
    tf_content = cb.text_frame
    tf_content.word_wrap = True
    
    # Split content by lines and add paragraphs
    lines = [line.strip() for line in content_text.split('\n') if line.strip()]
    
    if not lines:
        return

    first_line = True
    for line in lines:
        if first_line:
            p = tf_content.paragraphs[0]
            first_line = False
        else:
            p = tf_content.add_paragraph()
            
        p.text = line
        
        # Detect bold patterns like **xxx**:
        if line.startswith('**') and ':' in line:
            apply_font(p, size=18, bold=True, color=MAIN_TERRACOTTA)
        elif line.startswith('- ') or line.startswith('* '):
            p.level = 1
            apply_font(p, size=16, color=TEXT_DARK)
        else:
            apply_font(p, size=16, color=TEXT_DARK)

try:
    with open("PPT大纲_24页优化版.md", "r", encoding="utf-8") as f:
        md_text = f.read()
        
    # Split by slide marker (e.g., ### 第1页：)
    # Using regex to capture the slide blocks
    slide_blocks = re.split(r'### 第\d+页：', md_text)
    
    if len(slide_blocks) > 1:
        # Ignore block 0 (it's the document header)
        for i, block in enumerate(slide_blocks[1:]):
            lines = block.strip().split('\n')
            if not lines: continue
            
            # The first line is usually the slide title
            raw_title = lines[0].strip()
            # Clean up markdown bolding from title if present
            raw_title = raw_title.replace('**', '').replace('标题:', '').strip()
            
            content_lines = lines[1:]
            content_text = '\n'.join(content_lines)
            
            # Page 1 is usually cover
            if i == 0:
                # Extract subtitle and other info for cover
                subtitle_parts = []
                for cl in content_lines:
                    cl_clean = cl.replace('**', '').strip()
                    if cl_clean.startswith('副标题:') or cl_clean.startswith('演讲者:') or cl_clean.startswith('日期:'):
                        subtitle_parts.append(cl_clean.split(':', 1)[-1].strip())
                
                subtitle = '\n'.join(subtitle_parts) if subtitle_parts else "AI时代的高效工具"
                create_cover_slide("OpenClaw 使用指南", subtitle)
            else:
                create_content_slide(raw_title, content_text)
                
    prs.save("OpenClaw使用指南_24页完整莫兰迪风.pptx")
    print(f"Successfully generated full PPT with {len(prs.slides)} slides.")
except Exception as e:
    print(f"Error: {e}")
