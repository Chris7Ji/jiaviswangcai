#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
SECONDARY = RGBColor(0x2E, 0x5C, 0x8A)
ACCENT = RGBColor(0xF4, 0x7C, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)

def add_title_slide(title, subtitle, speaker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(28)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    box3 = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.8))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f'演讲者: {speaker}'
    p3.font.size = Pt(20)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f'第{num}部分'
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1.5))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY
    title_bg.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
    return slide

def add_two_column_slide(title, left_title, left_content, right_title, right_content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY
    title_bg.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左栏
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
    
    # 右栏
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
    return slide

def add_end_slide(title, contact):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = contact
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    return slide

print('开始生成35页PPT...')

# 第一部分：部署与安装（10页）
add_title_slide('OpenClaw 使用指南', '打造你的AI智能助理 - 从安装到精通', '季')
print('✅ 第1页：封面')

add_content_slide('关于我 & 课程目标', [
    '关于我（季）：',
    '• OpenClaw深度用户，华为昇腾生态建设者',
    '• AI技术爱好者，实战经验丰富',
    '',
    '今天你将学到：',
    '• OpenClaw核心概念与架构',
    '• 完整部署流程与最佳实践',
    '• 记忆系统与Skills深度应用',
    '• 6个真实案例实战演示',
    '• 从入门到精通的完整路径'
])
print('✅ 第2页：演讲者介绍')

add_content_slide('什么是OpenClaw？', [
    '定义：OpenClaw是一个AI智能体平台，让AI从"聊天"进化到"执行"',
    '',
    '核心区别：',
    '• ChatGPT/Claude：问答对话，知识回答',
    '• OpenClaw：任务执行，操作外部系统',
    '',
    '一句话总结：',
    'OpenClaw = 大模型大脑 + Skills手脚 + 记忆系统 + 多平台接入'
])
print('✅ 第3页：什么是OpenClaw')

add_content_slide('OpenClaw 核心架构详解', [
    '三层架构：',
    '',
    '1. Gateway（网关层）：消息路由、多平台接入、本地运行',
    '2. Agent（智能体层）：总指挥+8个专业Agent、模型管理',
    '3. Skills（技能层）：5400+精选技能、声明式配置',
    '',
    '类比理解：',
    '• 大脑 = 大语言模型（理解推理）',
    '• 神经系统 = Gateway（信息传递）',
    '• 手脚 = Skills（执行操作）',
    '• 记忆 = Memory系统（经验积累）'
])
print('✅ 第4页：核心架构')

add_content_slide('为什么选择OpenClaw？', [
    '5大核心优势：',
    '',
    '1. 完全开源：GitHub 297k+ Stars，代码透明',
    '2. 本地部署：数据不出本地，隐私安全可控',
    '3. 多平台支持：20+消息渠道，微信钉钉飞书等',
    '4. 主动执行：定时任务、自动监控，真正解放双手',
    '5. 生态丰富：13,729+注册技能，5,490+精选技能'
])
print('✅ 第5页：为什么选择OpenClaw')

add_content_slide('安装前准备清单', [
    '系统要求：',
    '• macOS 12+ / Linux / Windows 10+ (WSL2)',
    '• Node.js 18+ (推荐20 LTS)',
    '• RAM: 4GB+ (推荐8GB)',
    '',
    '必要账号：',
    '• 大模型API Key（Moonshot/Deepseek/Minimax）',
    '• 消息平台账号（飞书/钉钉等，可选）',
    '',
    '建议配置：VS Code、iTerm2、Git'
])
print('✅ 第6页：安装前准备')

add_content_slide('安装步骤详解', [
    'Step 1: 安装Node.js',
    '  brew install node@20',
    '',
    'Step 2: 安装OpenClaw CLI',
    '  npm install -g openclaw',
    '',
    'Step 3: 初始化配置',
    '  openclaw init && openclaw gateway start',
    '',
    'Step 4: 配置大模型',
    '  openclaw configure',
    '',
    'Step 5: 验证安装',
    '  openclaw status'
])
print('✅ 第7页：安装步骤')

add_content_slide('安装核心 Skills（信息获取类）', [
    '1. Clawsec - 安全防护',
    '   clawhub install clawsec',
    '   安装前自动安全扫描，三档评级',
    '',
    '2. Tavily Search - AI实时搜索',
    '   clawhub install tavily-search',
    '   解决知识截止问题，每月1000次免费',
    '',
    '3. Multi Search Engine - 多引擎聚合',
    '   clawhub install multi-search-engine',
    '   17个搜索引擎，无需API Key'
])
print('✅ 第8页：核心Skills（上）')

add_content_slide('安装核心 Skills（办公与进化类）', [
    '4. Self-Improving Agent - 自我进化',
    '   自动记录错误和经验，持续学习优化',
    '',
    '5. Proactive Agent - 主动服务',
    '   心跳机制，任务监控，主动提醒',
    '',
    '6. GitHub - 代码管理',
    '   自然语言操作GitHub',
    '',
    '7. Office-Automation - 办公助手',
    '   日程、邮件、文档、数据分析'
])
print('✅ 第9页：核心Skills（下）')

add_content_slide('安装后验证与常见问题', [
    '验证安装：',
    '• openclaw --version',
    '• openclaw gateway status',
    '• openclaw config list',
    '',
    '常见问题：',
    '• npm安装失败：使用国内镜像',
    '• Gateway启动失败：检查端口占用',
    '• API Key错误：重新配置',
    '• 技能安装失败：检查网络或手动安装'
])
print('✅ 第10页：安装验证')

# 第二部分：记忆系统与Skills（12页）
add_section_slide('二', '记忆系统与Skills深度应用')
print('✅ 第11页：第二部分分隔')

add_content_slide('记忆系统架构全景', [
    '为什么需要记忆系统？',
    '• 大模型每次对话都是"全新的"',
    '• 无法跨会话保持上下文',
    '• 无法积累经验和偏好',
    '',
    '三层记忆体系：',
    '1. 会话级记忆：当前对话，临时数据',
    '2. 项目级记忆：持久化存储，核心配置',
    '3. 技能级记忆：技能配置，行为规则'
])
print('✅ 第12页：记忆系统架构')

add_content_slide('核心记忆文件详解', [
    'SOUL.md - 身份定义：',
    '• 我是谁、核心原则、性格特点',
    '',
    'MEMORY.md - 核心记忆：',
    '• 核心决策、定时任务、技术配置',
    '',
    'AGENTS.md - Agent配置：',
    '• 8个Agent职责、任务派发规则',
    '',
    'USER.md - 用户画像：',
    '• 老板信息、技术兴趣、工作习惯',
    '',
    'TOOLS.md - 工具配置：',
    '• 邮箱配置、API密钥、常用路径'
])
print('✅ 第13页：核心记忆文件')

add_content_slide('学习记录系统实战', [
    '为什么记录学习？',
    '• AI会犯同样的错误',
    '• 用户纠正应该被记住',
    '• 成功经验应该复用',
    '',
    '三种学习记录：',
    '1. LEARNINGS.md - 成功经验、最佳实践',
    '2. ERRORS.md - 错误记录、故障排查',
    '3. FEATURE_REQUESTS.md - 功能需求、改进建议'
])
print('✅ 第14页：学习记录系统')

add_content_slide('记忆系统维护策略', [
    '每周回顾（周日）：',
    '• 检查新增学习记录',
    '• 识别重复模式（3次以上）',
    '• 提取核心学习点',
    '• 如记录>10条，执行归档',
    '',
    '每月归档（月末）：',
    '• 归档到archive/monthly/',
    '• 生成月度摘要',
    '• 清理活跃记录（≤10条）',
    '',
    '每季审查（季末）：',
    '• 审查MEMORY.md',
    '• 更新AGENTS.md',
    '• 季度总结归档'
])
print('✅ 第15页：维护策略')

add_content_slide('Skills深度应用 - 信息获取类', [
    'Tavily Search 高级用法：',
    '• 基础搜索：search.sh "关键词"',
    '• 深度搜索：search.sh "关键词" advanced 10',
    '• 时间过滤：搜索最近7天内容',
    '',
    'Multi Search Engine 高级用法：',
    '• 站内搜索：site:zhihu.com',
    '• 文件类型：filetype:pdf',
    '• 多引擎对比：中文/英文/技术内容'
])
print('✅ 第16页：Skills信息获取')

add_content_slide('Skills深度应用 - 自我进化类', [
    'Self-Improving Agent：',
    '• 启用Hook自动记录',
    '• 手动记录学习（纠正/错误/新方法）',
    '• 自动检索历史记录',
    '',
    'Proactive Agent：',
    '• 心跳机制：每15分钟自动唤醒',
    '• 任务监控：持续跟踪进行中任务',
    '• 主动提醒：不用问，主动推',
    '',
    '生成7个配置文件：SOUL.md、AGENTS.md等'
])
print('✅ 第17页：Skills自我进化')

add_content_slide('Skills深度应用 - 办公自动化类', [
    'Office-Automation 功能矩阵：',
    '• 日程管理：会议安排',
    '• 邮件处理：自动发送',
    '• 文档编辑：报告生成',
    '• 数据分析：Excel处理',
    '',
    'GitHub Skill 实战：',
    '• 搜索仓库：gh search repos',
    '• 管理Issues：gh issue list',
    '• PR审查：gh pr view'
])
print('✅ 第18页：Skills办公自动化')

add_content_slide('Skills开发入门', [
    'Skill目录结构：',
    '• SKILL.md（技能说明）',
    '• scripts/（可执行脚本）',
    '• assets/（资源文件）',
    '• _meta.json（元数据）',
    '',
    '开发流程：',
    '1. 创建目录',
    '2. 编写SKILL.md',
    '3. 编写执行脚本',
    '4. 测试调试',
    '5. 发布到ClawHub'
])
print('✅ 第19页：Skills开发')

add_content_slide('记忆检索优化实战', [
    '问题：Memory Search太慢（300秒超时）',
    '',
    '解决方案：',
    '1. 禁用慢速功能',
    '   memory_search.enabled = false',
    '',
    '2. 关键词索引系统（INDEX.md）',
    '   快速定位信息',
    '',
    '3. 分类存储学习记录',
    '   by-topic/email.md、cron.md等',
    '',
    '优化效果：5-300秒 → <1秒（300倍提升）'
])
print('✅ 第20页：记忆检索优化')

add_content_slide('故障排查与调试', [
    '常见问题诊断：',
    '',
    '响应慢：',
    '• 检查模型响应时间',
    '• 禁用慢速功能',
    '',
    '技能失效：',
    '• clawhub list 检查状态',
    '• 查看技能日志',
    '',
    '记忆丢失：',
    '• 检查记忆文件',
    '• 验证文件内容'
])
print('✅ 第21页：故障排查')

# 第三部分：实战案例（12页）
add_section_slide('三', '实战案例深度解析')
print('✅ 第22页：第三部分分隔')

add_content_slide('案例1 - 定时任务自动化', [
    '场景：每天早上自动获取AI新闻并发送邮件',
    '',
    '工作流程：',
    '06:30 定时触发 → Tavily搜索 → 整理Markdown',
    '→ 生成HTML → 发送邮件 → 飞书通知',
    '',
    '实际运行数据：',
    '• 已稳定运行3个月',
    '• 成功率99.5%',
    '• 平均耗时3分52秒',
    '',
    '扩展：OpenClaw新闻监控、健康科研监控'
])
print('✅ 第23页：案例1定时任务')

add_content_slide('案例2 - 邮件自动化系统', [
    '场景：多场景自动邮件发送',
    '',
    '邮件发送矩阵：',
    '• AI新闻：每天06:30 → QQ+华为邮箱',
    '• 健康报告：每天08:30 → QQ+华为邮箱',
    '• 股票预警：实时触发 → 老板',
    '• 任务完成：即时 → 老板',
    '',
    '技术架构：',
    '• auto_send_email.sh（AI新闻）',
    '• send_health_email.sh（健康报告）',
    '• SMTP配置（QQ邮箱）'
])
print('✅ 第24页：案例2邮件系统')

add_content_slide('案例3 - 语音交互系统', [
    '场景：语音输入 + 语音回复',
    '',
    '技术栈：',
    '• TTS引擎：Azure TTS',
    '• 音色：zh-CN-XiaoyiNeural',
    '• 语速：+30%',
    '• 输出：MP3',
    '',
    '使用场景：',
    '• 新闻播报、消息提醒',
    '• 内容朗读、多语言支持',
    '',
    '流程：文本 → Azure TTS → MP3 → 飞书'
])
print('✅ 第25页：案例3语音交互')

add_content_slide('案例4 - 智能内容生成', [
    '场景：自动撰写各类报告和文章',
    '',
    '内容生成矩阵：',
    '• AI新闻摘要：5分钟',
    '• 竞品分析报告：10分钟',
    '• 技术博客：15分钟',
    '• 微信公众号：20分钟',
    '',
    '实际案例：',
    '• 伊朗战争局势分析（3分钟）',
    '• OpenClaw技能总结（5分钟）'
])
print('✅ 第26页：案例4内容生成')

add_two_column_slide('案例5 - Agent团队协作',
    'Agent团队',
    [
        '总指挥(main)',
        '  任务识别与分发',
        '笔杆子(creator)',
        '  内容创作与编辑',
        '参谋(canmou)',
        '  深度研究与分析',
        '进化官(jinhua)',
        '  代码开发与优化',
        '交易官(jiaoyi)',
        '  股票监控与分析',
        '社区官(shequ)',
        '  社区运营与管理'
    ],
    '协作流程',
    [
        '1. 任务识别',
        '  分析用户需求',
        '  确定任务类型',
        '2. Agent派发',
        '  选择合适Agent',
        '  分配子任务',
        '3. 并行执行',
        '  多个Agent同时工作',
        '  提高效率',
        '4. 结果汇总',
        '  整合各Agent输出',
        '  统一汇报给老板'
    ]
)
print('✅ 第27页：案例5Agent协作')

add_content_slide('案例6 - 一键生成公众号文章 ⭐', [
    '场景：多Agent协作撰写微信公众号文章',
    '',
    '工作流程：',
    '1. 老板提出需求（主题、字数、风格）',
    '2. 参谋Agent深度调研（3分29秒）',
    '3. 笔杆子Agent撰写文章（50秒）',
    '4. 尝试API自动发布（遇到限制）',
    '5. 生成手工发布辅助工具',
    '',
    '成果：调研报告 + 2800字文章 + 发布辅助工具',
    '总耗时：4分19秒（传统方式需数小时）'
])
print('✅ 第28页：案例6公众号文章')

add_content_slide('案例6 - 技术亮点与经验总结', [
    '技术亮点：',
    '• 多Agent协作：参谋（左脑）+ 笔杆子（右脑）',
    '• 错误处理：API失败→第三方工具→手工辅助',
    '• 格式优化：Markdown → 微信友好格式',
    '',
    '经验总结：',
    '✓ 多Agent协作能大幅提升效率',
    '✓ 遇到限制时要有备选方案',
    '✓ 自动化+人工辅助是最佳实践',
    '✓ 持续优化工作流程'
])
print('✅ 第29页：案例6总结')

add_content_slide('高级技巧 - 性能优化', [
    '问题诊断：响应慢',
    '',
    '优化方案：',
    '1. 禁用慢速功能',
    '   memory_search.enabled = false',
    '',
    '2. 关键词索引系统（INDEX.md）',
    '',
    '3. 预加载核心文件',
    '   SOUL.md、MEMORY.md、INDEX.md',
    '',
    '优化效果：5-300秒 → <1秒（300倍提升）'
])
print('✅ 第30页：性能优化')

add_content_slide('高级技巧 - 安全与隐私', [
    '本地优先架构优势：',
    '',
    '数据安全：',
    '✓ 数据不出本地',
    '✓ 不依赖云端服务',
    '✓ 可完全离线运行',
    '',
    '隐私保护：',
    '✓ 聊天记录本地存储',
    '✓ API Key本地加密',
    '✓ 自主控制数据',
    '',
    '安全配置：文件权限、环境变量、定期备份'
])
print('✅ 第31页：安全与隐私')

add_content_slide('高级技巧 - 扩展与定制', [
    '自定义Agent：',
    '• 配置专属Agent',
    '• 指定模型和Skills',
    '',
    '自定义Skill开发：',
    '• 创建Skill目录',
    '• 编写SKILL.md',
    '• 编写执行脚本',
    '• 发布到ClawHub',
    '',
    '集成外部服务：',
    '• 企业微信、钉钉、Slack'
])
print('✅ 第32页：扩展与定制')

add_content_slide('故障排查完整指南', [
    '系统化排查流程：',
    '',
    'Step 1: 收集信息',
    '• openclaw --version',
    '• openclaw status',
    '• 查看日志',
    '',
    'Step 2: 定位问题',
    '• 网络连接测试',
    '• API Key验证',
    '• 技能测试',
    '',
    'Step 3: 解决问题',
    '• 重启Gateway',
    '• 重新配置',
    '• 更新到最新版'
])
print('✅ 第33页：故障排查指南')

add_content_slide('价值总结与最佳实践', [
    'OpenClaw核心价值：',
    '• 效率提升：自动化任务，效率提升10倍+',
    '• 成本降低：开源免费，本地部署',
    '• 能力扩展：5400+技能生态，持续进化',
    '',
    '最佳实践：',
    '✓ DO：定期备份、记录经验、使用版本控制',
    '✗ DON\'T：硬编码敏感信息、忽视安全扫描'
])
print('✅ 第34页：价值总结')

add_end_slide('谢谢聆听！', 'Q&A 欢迎交流')
print('✅ 第35页：结束页')

# 保存
output_path = '/Users/jiyingguo/.openclaw/workspace/OpenClaw使用指南_35页完整版.pptx'
prs.save(output_path)
print(f'\n🎉 35页PPT生成完成！')
print(f'📁 保存路径: {output_path}')
print(f'📊 总页数: {len(prs.slides)} 页')