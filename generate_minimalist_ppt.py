#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
极简商务风PPT生成器
风格特征：
- 主色调：深炭灰 #2D3436
- 辅助色：浅灰 #F5F6FA 背景，白色 #FFFFFF 内容区
- 点缀色：珊瑚橙 #FF6B6B
- 字体：微软雅黑
- 留白：内容区域占比不超过60%
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 极简商务风色彩系统
COLORS = {
    'primary': RGBColor(45, 52, 54),      # 深炭灰 #2D3436
    'secondary': RGBColor(245, 246, 250),  # 浅灰 #F5F6FA
    'white': RGBColor(255, 255, 255),      # 白色 #FFFFFF
    'accent': RGBColor(255, 107, 107),     # 珊瑚橙 #FF6B6B
    'text': RGBColor(60, 60, 60),          # 正文灰
    'light_text': RGBColor(120, 120, 120), # 浅灰文字
}

def create_minimalist_ppt(slides_data, output_path):
    """创建极简商务风PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides_data:
        layout_type = slide_data.get('layout', 'content')
        
        if layout_type == 'cover':
            add_cover_slide(prs, slide_data)
        elif layout_type == 'section':
            add_section_slide(prs, slide_data)
        elif layout_type == 'two-column':
            add_two_column_slide(prs, slide_data)
        elif layout_type == 'three-column':
            add_three_column_slide(prs, slide_data)
        elif layout_type == 'timeline':
            add_timeline_slide(prs, slide_data)
        elif layout_type == 'matrix':
            add_matrix_slide(prs, slide_data)
        elif layout_type == 'table':
            add_table_slide(prs, slide_data)
        else:
            add_content_slide(prs, slide_data)
    
    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")

def add_cover_slide(prs, data):
    """封面页 - 极简设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景色块 - 左侧深炭灰
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    
    # 标题 - 白色，左侧
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题
    if data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(4), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['accent']
        p.alignment = PP_ALIGN.LEFT
    
    # 演讲者信息 - 右侧
    info_box = slide.shapes.add_textbox(Inches(6), Inches(5.5), Inches(6), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    if data.get('speaker'):
        p = tf.paragraphs[0]
        p.text = f"演讲者: {data['speaker']}"
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['light_text']
    
    if data.get('date'):
        p = tf.add_paragraph()
        p.text = data['date']
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['light_text']
        p.space_before = Pt(8)

def add_section_slide(prs, data):
    """章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['secondary']
    background.line.fill.background()
    
    # 章节编号
    if data.get('number'):
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{data['number']}" if data['number'] < 10 else str(data['number'])
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(10), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 副标题
    if data.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(10), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['light_text']

def add_content_slide(prs, data):
    """标准内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 标题下划线 - 珊瑚橙点缀
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    content = data.get('content', [])
    if isinstance(content, str):
        content = [content]
    
    for i, item in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 处理加粗文本（冒号前加粗）
        if '：' in item:
            parts = item.split('：', 1)
            run1 = p.add_run()
            run1.text = parts[0] + '：'
            run1.font.name = 'Microsoft YaHei'
            run1.font.size = Pt(14)
            run1.font.bold = True
            run1.font.color.rgb = COLORS['primary']
            
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.name = 'Microsoft YaHei'
            run2.font.size = Pt(14)
            run2.font.color.rgb = COLORS['text']
        else:
            p.text = item
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text']
        
        p.space_after = Pt(12)
        p.line_spacing = 1.5

def add_two_column_slide(prs, data):
    """双栏布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 标题下划线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    columns = data.get('columns', [])
    
    # 左栏
    if len(columns) > 0:
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        # 栏目标题
        p = tf.paragraphs[0]
        p.text = columns[0].get('title', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        p.space_after = Pt(12)
        
        # 内容
        for item in columns[0].get('content', []):
            p = tf.add_paragraph()
            p.text = item
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(8)
            p.line_spacing = 1.4
    
    # 右栏
    if len(columns) > 1:
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(5.2))
        tf = right_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = columns[1].get('title', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        p.space_after = Pt(12)
        
        for item in columns[1].get('content', []):
            p = tf.add_paragraph()
            p.text = item
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(8)
            p.line_spacing = 1.4

def add_three_column_slide(prs, data):
    """三栏布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    columns = data.get('columns', [])
    col_width = 3.8
    col_spacing = 0.4
    
    for i, col in enumerate(columns[:3]):
        left = Inches(0.8 + i * (col_width + col_spacing))
        col_box = slide.shapes.add_textbox(left, Inches(1.6), Inches(col_width), Inches(5.2))
        tf = col_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = col.get('title', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        p.space_after = Pt(10)
        
        for item in col.get('content', []):
            p = tf.add_paragraph()
            p.text = item
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(6)
            p.line_spacing = 1.3

def add_timeline_slide(prs, data):
    """时间线布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    steps = data.get('steps', [])
    step_width = 11.5 / len(steps) if steps else 3
    
    for i, step in enumerate(steps):
        left = Inches(0.8 + i * step_width)
        
        # 时间节点圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(step_width/2 - 0.15), Inches(2), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLORS['accent']
        dot.line.fill.background()
        
        # 步骤标题
        step_box = slide.shapes.add_textbox(left, Inches(2.5), Inches(step_width - 0.2), Inches(0.8))
        tf = step_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step.get('title', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 步骤描述
        desc_box = slide.shapes.add_textbox(left, Inches(3.4), Inches(step_width - 0.2), Inches(3))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step.get('desc', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER

def add_matrix_slide(prs, data):
    """矩阵布局 (2x2)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    quadrants = data.get('quadrants', [])
    positions = [
        (Inches(0.8), Inches(1.6)),   # 左上
        (Inches(6.8), Inches(1.6)),   # 右上
        (Inches(0.8), Inches(4.3)),   # 左下
        (Inches(6.8), Inches(4.3)),   # 右下
    ]
    
    for i, quad in enumerate(quadrants[:4]):
        left, top = positions[i]
        
        # 象限背景
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.5), Inches(2.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['secondary']
        bg.line.fill.background()
        
        # 象限标题
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.1), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = quad.get('title', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        
        # 象限描述
        desc_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9), Inches(5.1), Inches(1.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quad.get('desc', '')
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text']

def add_table_slide(prs, data):
    """表格布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    rows = data.get('rows', [])
    if not rows:
        return
    
    headers = rows[0]
    data_rows = rows[1:]
    
    # 表格位置
    table_left = Inches(0.8)
    table_top = Inches(1.6)
    table_width = Inches(11.5)
    table_height = Inches(0.5 + len(data_rows) * 0.5)
    
    # 创建表格
    table = slide.shapes.add_table(len(data_rows) + 1, len(headers), table_left, table_top, table_width, table_height).table
    
    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = COLORS['white']
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
    
    # 设置数据行
    for row_idx, row_data in enumerate(data_rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = COLORS['text']
            
            # 隔行变色
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['secondary']

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python generate_minimalist_ppt.py <input.json> <output.pptx>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    with open(input_file, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    create_minimalist_ppt(slides_data, output_file)
