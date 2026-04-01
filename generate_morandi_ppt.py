from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

# Colors
MAIN_TERRACOTTA = RGBColor(165, 129, 105) # #A58169
MAIN_SAGE = RGBColor(111, 122, 97) # #6F7A61
BG_MORANDI_GRAY = RGBColor(240, 235, 229) # #F0EBE5
CONTENT_WARM_WHITE = RGBColor(253, 251, 246) # #FDFBF6
ACCENT_ROSE = RGBColor(201, 160, 155) # #C9A09B
TEXT_DARK = RGBColor(80, 80, 80)

def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_MORANDI_GRAY

def create_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    
    title_box = slide.shapes.title
    title_box.text = title
    for p in title_box.text_frame.paragraphs:
        p.font.name = 'Source Han Sans CN'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = MAIN_TERRACOTTA
        p.alignment = PP_ALIGN.CENTER
        
    sub_box = slide.placeholders[1]
    sub_box.text = subtitle
    for p in sub_box.text_frame.paragraphs:
        p.font.name = 'Source Han Sans CN'
        p.font.size = Pt(24)
        p.font.color.rgb = MAIN_SAGE
        p.alignment = PP_ALIGN.CENTER
    return slide

def create_content_slide(title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_bg(slide)
    
    # White content box background
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.5)
    rect = slide.shapes.add_shape(1, left, top, width, height) # 1 = rectangle
    rect.fill.solid()
    rect.fill.fore_color.rgb = CONTENT_WARM_WHITE
    rect.line.fill.background()
    
    title_box = slide.shapes.title
    title_box.text = title
    for p in title_box.text_frame.paragraphs:
        p.font.name = 'Source Han Sans CN'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = MAIN_SAGE
        p.alignment = PP_ALIGN.LEFT
        
    content_box = slide.placeholders[1]
    content_box.left = Inches(1.0)
    content_box.top = Inches(2.0)
    content_box.width = Inches(8.0)
    content_box.height = Inches(4.5)
    
    tf = content_box.text_frame
    tf.text = content_list[0] if content_list else ""
    for p in tf.paragraphs:
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK
    
    for item in content_list[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK

create_title_slide("OpenClaw 使用指南", "打造你的AI智能助理\n季英国 | 2026年3月")
create_content_slide("关于课程目标", ["今天你将学到：", "• OpenClaw核心概念与架构", "• 完整部署流程与最佳实践", "• 记忆系统与Skills深度应用", "• 6个真实案例演示"])

prs.save("OpenClaw使用指南_极简莫兰迪风.pptx")
print("PPT saved.")
